from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typo_check import check_spelling
import uuid

EMU_PER_CM = 360000
EMU_PER_PT = 12700

def emu_to_cm(emu_value):
    return emu_value / EMU_PER_CM

def emu_to_pt(emu_value):
    return emu_value / EMU_PER_PT if emu_value is not None else None

def spellcheck_batch_with_markers(texts):
    """
    texts: List[str]
    return: List[str or None]  # None이면 해당 항목만 폴백 권장
    """
    if not texts:
        return []

    # 1) 각 텍스트에 유니크 마커 부착
    segments = []  # [(start, end)]
    pieces = []
    for t in texts:
        uid = uuid.uuid4().hex
        start = f"<<<S_{uid}>>>"
        end   = f"<<<E_{uid}>>>"
        segments.append((start, end))
        # 줄바꿈을 일부 교정기가 만지더라도 마커로 경계를 잡음
        pieces.append(f"{start}{t}{end}")

    big_payload = "\n".join(pieces)

    # 2) 한 번만 호출
    corrected_big = check_spelling(big_payload)

    # 3) 마커로 다시 분해
    results = []
    cursor = 0
    for start, end in segments:
        s = corrected_big.find(start, cursor)
        if s == -1:
            results.append(None)
            continue
        s += len(start)
        e = corrected_big.find(end, s)
        if e == -1:
            results.append(None)
            cursor = s
            continue
        results.append(corrected_big[s:e])
        cursor = e + len(end)

    return results

# ===== 메인 로직 =====
prs = Presentation("input/ppt/example3.pptx")

# (A) 먼저 모든 텍스트 수집
all_texts = []
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    t = run.text.strip()
                    if t:
                        all_texts.append(t)

# (B) 한 번만 맞춤법 검사 (마커 방식)
corrected_results = spellcheck_batch_with_markers(all_texts)

# (C) 출력 루프 (fallback 제거)
result_idx = 0
for slide_idx, slide in enumerate(prs.slides, start=1):
    print(f"\n=== Slide {slide_idx} ===")

    for shape_idx, shape in enumerate(slide.shapes, start=1):
        print(f"\n[Shape {shape_idx}]")
        print(f"  Shape type: {shape.shape_type}")

        left_cm = emu_to_cm(shape.left)
        top_cm = emu_to_cm(shape.top)
        width_cm = emu_to_cm(shape.width)
        height_cm = emu_to_cm(shape.height)

        print(f"  위치 (left, top): {left_cm:.2f} cm, {top_cm:.2f} cm")
        print(f"  크기 (width x height): {width_cm:.2f} cm x {height_cm:.2f} cm")

        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    text = run.text
                    if text.strip():
                        # 배치 결과만 사용
                        corrected_text = corrected_results[result_idx] \
                            if result_idx < len(corrected_results) else text

                        print(f"  원본 텍스트: {text}")
                        print(f"  교정된 텍스트: {corrected_text}")
                        result_idx += 1

                    font = run.font.name
                    size_pt = emu_to_pt(run.font.size)
                    bold = run.font.bold
                    italic = run.font.italic
                    print(f"    폰트: {font}, 크기: {size_pt} pt, Bold: {bold}, Italic: {italic}")
