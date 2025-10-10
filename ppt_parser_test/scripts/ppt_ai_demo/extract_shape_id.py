from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

prs = Presentation("before.pptx")

for slide_idx, slide in enumerate(prs.slides):
    print(f"--- Slide {slide_idx} ---")
    for shape in slide.shapes:
        shape_type = MSO_SHAPE_TYPE(shape.shape_type).name
        has_text = getattr(shape, "has_text_frame", False)

        # 기본 정보
        info = f"shape_id={shape.shape_id}, type={shape_type}, has_text={has_text}"

        # 텍스트가 있으면 내용도 출력
        if has_text and shape.text_frame.text.strip():
            info += f", text='{shape.text_frame.text.strip()}'"

        # 위치와 크기 출력
        info += f", left={shape.left}, top={shape.top}, width={shape.width}, height={shape.height}"

        # 도형이면 fill color 출력 (단색 기준, 투명도 제외)
        if shape_type in ["AUTO_SHAPE", "FREEFORM", "TEXT_BOX"]:
            try:
                fill = shape.fill
                if fill.type is not None and hasattr(fill.fore_color, "rgb") and fill.fore_color.rgb is not None:
                    rgb = fill.fore_color.rgb
                    info += f", fill_color=[{rgb[0]}, {rgb[1]}, {rgb[2]}]"
                else:
                    info += ", fill_color=None"
            except Exception:
                info += ", fill_color=None"

        print(info)
