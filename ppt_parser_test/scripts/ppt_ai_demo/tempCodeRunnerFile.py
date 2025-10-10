from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageEnhance
import io
import json
from lxml import etree
from pptx.oxml.ns import qn

# JSON 불러오기
with open("scripts/ppt_ai_demo/improve.json", "r", encoding="utf-8") as f:
    instructions = json.load(f)

# PPT 불러오기
prs = Presentation("before.pptx")

# 도형 투명도 설정
def set_fill_transparency(shape, transparency):
    """
    shape: pptx shape
    transparency: 0~1
    """
    # Shape: solid fill
    if hasattr(shape, "fill"):
        shape.fill.solid()
        sp = shape._element
        solidFill = sp.find(".//a:solidFill", namespaces=sp.nsmap)
        if solidFill is None:
            solidFill = etree.SubElement(sp, "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")

        srgbClr = solidFill.find(".//a:srgbClr", namespaces=sp.nsmap)
        if srgbClr is None:
            srgbClr = etree.SubElement(solidFill, "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")

        # 기존 RGB 가져오기
        rgb = shape.fill.fore_color.rgb
        srgbClr.set("val", "%02X%02X%02X" % (rgb[0], rgb[1], rgb[2]))

        alpha = int(transparency * 100000)  # pptx 단위
        alpha_elem = srgbClr.find(".//a:alpha", namespaces=sp.nsmap)
        if alpha_elem is None:
            alpha_elem = etree.SubElement(srgbClr, "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha")
        alpha_elem.set("val", str(alpha))


def apply_transparency_to_image(img_path, transparency):
    """
    img_path: 원본 이미지 경로
    transparency: 0~1 (0=투명, 1=불투명)
    """
    im = Image.open(img_path).convert("RGBA")
    alpha = im.split()[3]  # 알파 채널
    alpha = alpha.point(lambda p: int(p * transparency))
    im.putalpha(alpha)
    # 메모리 버퍼에 저장
    img_byte_arr = io.BytesIO()
    im.save(img_byte_arr, format='PNG')
    img_byte_arr.seek(0)
    return img_byte_arr

# JSON 기반으로 슬라이드/도형/이미지 수정
for slide_idx, slide_instr in enumerate(instructions.get("slides", [])):
    if slide_idx >= len(prs.slides):
        continue
    slide = prs.slides[slide_idx]

    for shape_instr in slide_instr.get("shapes", []):
        # ID로 shape 찾기
        shape = None
        for sh in slide.shapes:
            if sh.shape_id == shape_instr.get("id"):
                shape = sh
                break
        if shape is None:
            continue

        # ===== 텍스트 수정 =====
        if shape_instr.get("type") == "text" and getattr(shape, "has_text_frame", False):
            tf = shape.text_frame

            if "text" in shape_instr:
                tf.clear()  # 기존 텍스트/스타일 전부 지우기
                for line in shape_instr["text"].split("\n"):
                    p = tf.add_paragraph()
                    r = p.add_run()
                    r.text = line

                    if "font_size" in shape_instr:
                        r.font.size = Pt(shape_instr["font_size"])
                    if "color" in shape_instr:
                        r.font.color.rgb = RGBColor(*shape_instr["color"])
                    if "bold" in shape_instr:
                        r.font.bold = shape_instr["bold"]
                    if "font" in shape_instr:
                        r.font.name = shape_instr["font"]

                    if "line_spacing" in shape_instr:
                        p.line_spacing = shape_instr["line_spacing"]
                    # 글머리표 처리
                    if not shape_instr.get("hasBullet", True):
                        pPr = p._element.get_or_add_pPr()
                        buNone = etree.Element(qn('a:buNone'))
                        pPr.append(buNone)
            # 위치 조정
            if "left" in shape_instr:
                shape.left = shape_instr["left"]
            if "top" in shape_instr:
                shape.top = shape_instr["top"]

        # ===== 이미지 수정 =====
        elif shape_instr.get("type") == "image" and shape.shape_type == 13:  # PICTURE
            for attr in ["width", "height", "left", "top"]:
                if attr in shape_instr:
                    setattr(shape, attr, shape_instr[attr])
            # 투명도 적용
            if "transperency" in shape_instr:
                # 기존 이미지 파일 경로 가져오기
                img_part = shape.image
                img_bytes = img_part.blob
                img_stream = io.BytesIO(img_bytes)
                # Pillow로 투명도 적용
                transparent_img = apply_transparency_to_image(img_stream, shape_instr["transperency"])
                # 기존 이미지를 대체
                shape._element.getparent().remove(shape._element)  # 기존 이미지 제거
                slide.shapes.add_picture(transparent_img, shape.left, shape.top, width=shape.width, height=shape.height)

        # ===== 도형 수정 =====
        elif shape_instr.get("type") == "shape" and hasattr(shape, "fill"):
            # 색상 변경
            if "fill_color" in shape_instr:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*shape_instr["fill_color"])
            # 위치와 크기
            for attr in ["width", "height", "left", "top"]:
                if attr in shape_instr:
                    setattr(shape, attr, shape_instr[attr])
            # 투명도 적용
            if "transperency" in shape_instr:
                set_fill_transparency(shape, shape_instr["transperency"])

# PPT 저장
prs.save("after.pptx")
