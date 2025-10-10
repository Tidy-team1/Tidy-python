from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import time

EMU_PER_CM = 360000   # 1cm = 360,000 EMU
EMU_PER_PT = 12700    # 1pt = 12,700 EMU

def emu_to_cm(emu_value):
    return emu_value / EMU_PER_CM

def emu_to_pt(emu_value):
    return emu_value / EMU_PER_PT if emu_value is not None else None

start_time = time.time()  # 시작 시간 기록

# 테스트할 PPTX 파일 열기
prs = Presentation("input/ppt/example3.pptx")

for slide_idx, slide in enumerate(prs.slides, start=1):
    print(f"\n=== Slide {slide_idx} ===")

    for shape_idx, shape in enumerate(slide.shapes, start=1):
        print(f"\n[Shape {shape_idx}]")
        print(f"  Shape type: {shape.shape_type}")  # 도형 종류 (텍스트박스, 사각형, 이미지 등)

        left_cm = emu_to_cm(shape.left)
        top_cm = emu_to_cm(shape.top)
        width_cm = emu_to_cm(shape.width)
        height_cm = emu_to_cm(shape.height)

        print(f"  위치 (left, top): {left_cm:.2f} cm, {top_cm:.2f} cm")
        print(f"  크기 (width x height): {width_cm:.2f} cm x {height_cm:.2f} cm")

        # 텍스트가 있는 도형 처리
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    text = run.text
                    font = run.font.name
                    size_pt = emu_to_pt(run.font.size)  # pt 단위 변환
                    bold = run.font.bold
                    italic = run.font.italic
                    print(f"  텍스트: {text}")
                    print(f"    폰트: {font}, 크기: {size_pt} pt, Bold: {bold}, Italic: {italic}")

        # 도형(예: 사각형, 원 등) 처리
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            print("  → 도형 객체 (AutoShape)")
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print("  → 이미지 객체 (Picture)")

end_time = time.time()    # 끝난 시간 기록
elapsed_time = end_time - start_time  # 경과 시간 계산
print(f"검사 소요 시간: {elapsed_time:.2f}초")  # 초 단위 출력