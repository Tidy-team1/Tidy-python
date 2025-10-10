from pptx import Presentation
from pptx.util import Inches
from PIL import Image

# 이미지 불러오기
img = Image.open("input/image/new_img.png").convert("RGBA")

# 알파값 조절 (예: 50%)
alpha = img.split()[3]
alpha = alpha.point(lambda p: p * 0.5)
img.putalpha(alpha)

# 수정된 이미지 저장
img.save("output/image/transparent.png")

# PPTX에 삽입
prs = Presentation("input/ppt/example.pptx")
slide = prs.slides[0]
slide.shapes.add_picture("output/image/transparent.png", Inches(2), Inches(2), Inches(3), Inches(2))
prs.save("output/ppt/example_transparent_final.pptx")
