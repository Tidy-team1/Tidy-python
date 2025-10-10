from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 재귀적으로 도형 탐색
def extract_shapes(shapes, depth=0):
    results = []
    for shape in shapes:
        indent = "  " * depth  # 그룹 계층 확인용 들여쓰기

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # 그룹 정보
            results.append(f"{indent}[Group] pos=({shape.left}, {shape.top}), size=({shape.width}, {shape.height})")
            # 그룹 내부 순회 (재귀 호출)
            results.extend(extract_shapes(shape.shapes, depth + 1))

        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # 이미지 정보
            results.append(f"{indent}[Picture] size=({shape.width}, {shape.height})")
            # 이미지 저장 (선택 사항)
            image = shape.image
            with open(f"extracted_image_{id(shape)}.{image.ext}", "wb") as f:
                f.write(image.blob)

        elif shape.has_text_frame:
            # 텍스트 상자/도형 안의 텍스트
            text = shape.text.strip()
            if text:
                results.append(f"{indent}[Text] \"{text}\"")

        else:
            # 그 외 기본 도형
            results.append(f"{indent}[Shape] type={shape.shape_type}")

    return results


if __name__ == "__main__":
    prs = Presentation("input/ppt/example3.pptx")

    for i, slide in enumerate(prs.slides, start=1):
        print(f"\n=== Slide {i} ===")
        extracted = extract_shapes(slide.shapes)
        for item in extracted:
            print(item)
