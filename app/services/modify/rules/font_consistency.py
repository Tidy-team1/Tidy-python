# app/services/modify/rules/font_consistency.py

from pptx.util import Pt
from app.core.logger import logger

def apply_font_consistency(slide, item, details):
    if item.shapeId is None:
        logger.warning("[font_consistency] shapeId is None → skip")
        return

    current_size = details.get("currentSize")
    recommended = details.get("recommendedSize")

    if not current_size or not recommended:
        logger.warning(
            f"[font_consistency] currentSize/recommendedSize missing for shapeId={item.shapeId}"
        )
        return

    updated = False

    for shape in slide.shapes:
        if getattr(shape, "shape_id", None) != item.shapeId:
            continue
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font:
                    run.font.size = Pt(recommended)

        updated = True
        logger.info(
            f"[font_consistency] shapeId={item.shapeId} font {current_size} → {recommended}"
        )

    if not updated:
        logger.warning(f"[font_consistency] No shape updated for shapeId={item.shapeId}")
