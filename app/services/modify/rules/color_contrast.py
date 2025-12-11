# app/services/modify/rules/color_contrast.py

from app.core.logger import logger
from pptx.dml.color import RGBColor
from app.utils.color_utils import hex_to_rgb_int


def apply_color_contrast(slide, item, details):
    """
    details 예:
      {
        "current": "#F0E4EE",
        "contrast": 1.23,
        "required": 3.0,
        "recommended": ["#6E567A", "#1A191A", "#FFFFFF"]
      }
    recommended[0] 색상으로 텍스트 색을 변경한다.
    """

    shape_id = item.shapeId
    if shape_id is None:
        logger.warning("[color_contrast] shapeId is None → skip")
        return

    recommended = details.get("recommended")
    if not recommended or not isinstance(recommended, list):
        logger.warning(f"[color_contrast] No recommended colors for shapeId={shape_id}")
        return

    new_color_hex = recommended[0]
    r, g, b = hex_to_rgb_int(new_color_hex)

    updated = False

    # shape 찾기
    for shape in slide.shapes:
        if getattr(shape, "shape_id", None) != shape_id:
            continue
        if not shape.has_text_frame:
            continue

        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font:
                    run.font.color.rgb = RGBColor(r, g, b)

        updated = True
        logger.info(
            f"[color_contrast] shapeId={shape_id} color changed → {new_color_hex}"
        )

    if not updated:
        logger.warning(f"[color_contrast] No shape updated for shapeId={shape_id}")
