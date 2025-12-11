import cv2
import numpy as np
import io
from pptx.enum.shapes import MSO_SHAPE_TYPE
from app.core.logger import logger
from app.utils.ppt_utils import emu_to_px

def apply_image_contrast(slide, item, details):
    """
    이미지의 대비(Contrast)나 투명도(Alpha)를 조절하여 다시 삽입합니다.
    
    details 구조 예시:
    {
        "current": "Image",
        "contrast": 20.5,
        "alpha_score": 100.0,
        "recommended": {
            "contrast_increase": 15.0, # 부족한 값
            "target_contrast": 40.0,   # 목표 값
            "alpha_increase": 50.0,
            "target_alpha": 180.0
        }
    }
    """
    shape_id = item.shapeId
    recommended = details.get("recommended")

    if not recommended:
        logger.warning(f"[image_contrast] No adjustments found for shapeId={shape_id}")
        return

    # 1. 대상 Shape 찾기
    target_shape = None
    for shape in slide.shapes:
        if getattr(shape, "shape_id", None) == shape_id:
            target_shape = shape
            break
    
    if not target_shape:
        logger.warning(f"[image_contrast] Shape not found: shapeId={shape_id}")
        return

    if target_shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        logger.warning(f"[image_contrast] Shape is not a picture: shapeId={shape_id}")
        return

    # 2. 이미지 추출 및 OpenCV 변환
    try:
        image_blob = target_shape.image.blob
        nparr = np.frombuffer(image_blob, np.uint8)
        img = cv2.imdecode(nparr, cv2.IMREAD_UNCHANGED) # Alpha 채널 포함 로드
    except Exception as e:
        logger.error(f"[image_contrast] Failed to load image blob: {e}")
        return

    # 3. 이미지 보정 (대비 및 투명도)
    processed_img = _process_image_opencv(img, details, recommended)
    
    if processed_img is None:
        logger.warning("[image_contrast] Image processing returned None")
        return

    # 4. 보정된 이미지를 바이트로 변환
    # 원본 포맷 유지가 좋으나, 투명도 보존을 위해 PNG 권장
    success, encoded_img = cv2.imencode(".png", processed_img)
    if not success:
        logger.error("[image_contrast] Failed to encode processed image")
        return
    
    new_image_stream = io.BytesIO(encoded_img.tobytes())

    # 5. 슬라이드에서 이미지 교체 (기존 위치/크기 보존)
    _replace_picture_in_slide(slide, target_shape, new_image_stream)
    
    logger.info(f"[image_contrast] applied successfully for shapeId={shape_id}")

def _process_image_opencv(img, details, adjustments):
    """
    [개선된 알고리즘]
    감마 보정(Gamma Correction)을 사용하여
    흰색(255)은 유지하되, 흐릿한 회색 글씨를 강제로 진하게(검정에 가깝게)
    """
    res = img.copy()
    
    # 1. Alpha 채널 분리
    has_alpha = False
    if len(res.shape) == 3 and res.shape[2] == 4:
        has_alpha = True
        b, g, r, a = cv2.split(res)
        bgr = cv2.merge([b, g, r])
    else:
        bgr = res

    # 2. 감마 보정 적용 (핵심 로직)
    # Gamma < 1.0: 이미지가 어두워짐 (흐릿한 글씨가 진해짐)
    # Gamma > 1.0: 이미지가 밝아짐
    # 여기서는 글씨를 선명하게 하기 위해 0.5 정도의 강한 값을 적용
    gamma_value = 0.5 
    
    # 감마 테이블 생성 (Lookup Table)
    invGamma = 1.0 / gamma_value
    table = np.array([((i / 255.0) ** invGamma) * 255 for i in np.arange(0, 256)]).astype("uint8")
    
    # LUT 적용
    enhanced_bgr = cv2.LUT(bgr, table)

    # (선택 사항) Min-Max 정규화 추가 적용: 가장 어두운 픽셀을 완전 검정(0)으로 만듦
    enhanced_bgr = cv2.normalize(enhanced_bgr, None, 0, 255, cv2.NORM_MINMAX)

    # 3. Alpha 채널 다시 합치기 & 투명도 보정
    if has_alpha:
        if "target_alpha" in adjustments:
            target_a = adjustments["target_alpha"]
            current_a = details.get("alpha_score", 255.0)
            if current_a < target_a and current_a > 0:
                alpha_scale = target_a / current_a
                alpha_scale = min(alpha_scale, 2.5)
                a = np.clip(a * alpha_scale, 0, 255).astype(np.uint8)
        
        res = cv2.merge([enhanced_bgr[:,:,0], enhanced_bgr[:,:,1], enhanced_bgr[:,:,2], a])
    else:
        res = enhanced_bgr

    return res

def _replace_picture_in_slide(slide, old_shape, new_image_stream):
    """
    기존 shape의 위치/크기를 가져와서 새 이미지를 삽입하고 기존 shape을 삭제합니다.
    (python-pptx는 기존 이미지 blob의 직접 교체가 어렵기 때문에 삭제 후 재삽입 방식 사용)
    """
    # 1. 기존 속성 저장
    left = old_shape.left
    top = old_shape.top
    width = old_shape.width
    height = old_shape.height
    name = old_shape.name
    
    # 2. 새 이미지 추가
    new_pic = slide.shapes.add_picture(new_image_stream, left, top, width, height)
    new_pic.name = name # 이름 유지 (선택사항)

    # 3. 기존 이미지 삭제
    # python-pptx에는 delete() 메서드가 없으므로 element 트리에서 제거해야 함
    sp = old_shape._element
    sp.getparent().remove(sp)