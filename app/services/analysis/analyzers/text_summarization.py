# app/services/analysis/analyzers/text_summarization.py
from pptx import Presentation
from app.services.analysis.base_analyzer import BaseAnalyzer
from app.models.analysis_dto import IssueElement, IssueResult, SlideIssueResult
from app.utils.text_summarizer import TextSummarizer
from app.utils.ppt_utils import emu_to_px

class TextSummarizationAnalyzer(BaseAnalyzer):
    analyzer_type = "text_summarization"

    def __init__(self):
        self.summarizer = TextSummarizer(model="gpt-4.1") # 모델명은 환경에 맞게
        # 개별 텍스트 박스 기준 임계값 (조금 더 엄격하게 잡아도 됨)
        self.min_lines = 4 
        self.min_chars = 40 

    def analyze(self, prs: Presentation) -> list[SlideIssueResult]:
        results = []
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_issues = []
            
            # [변경] 슬라이드 전체 텍스트가 아니라, 각 Shape(요소)별로 순회
            for i, shape in enumerate(slide.shapes):
                
                # 1. 텍스트 프레임이 있고, 텍스트가 존재하는지 확인
                if not shape.has_text_frame or not shape.text.strip():
                    continue

                raw_text = shape.text.strip()
                lines = raw_text.splitlines()

                # 2. [임계값 체크] 해당 박스의 텍스트가 요약이 필요한 수준인가?
                if len(lines) >= self.min_lines or len(raw_text) >= self.min_chars:
                    
                    # 3. 요약 수행 (해당 텍스트 박스 내용만)
                    summary_result = self.summarizer.summarize(raw_text)
                    
                    if summary_result and "summary_bullets" in summary_result:
                        recommended = summary_result["summary_bullets"]
                    else:
                        # 요약 실패 시 스킵하거나 에러 메시지
                        continue 

                    # 4. 이슈 등록 (Shape ID 필수)
                    issue_elem = IssueElement(
                        shapeId=shape.shape_id,
                        elementIndex=i,
                        bboxLeft=emu_to_px(shape.left),
                        bboxTop=emu_to_px(shape.top),
                        bboxWidth=emu_to_px(shape.width),
                        bboxHeight=emu_to_px(shape.height),
                        text=raw_text,
                        elementType="text_box"
                    )

                    slide_issues.append(
                        IssueResult(
                            type=self.analyzer_type,
                            message="텍스트 박스의 내용이 너무 길어 가독성이 떨어집니다.",
                            element=issue_elem,
                            details={
                                "current": raw_text,
                                "recommend": recommended # 리스트 형태 ["요약1", "요약2"]
                            }
                        )
                    )

            if slide_issues:
                results.append(SlideIssueResult(slide=slide_idx, issues=slide_issues))
        
        return results