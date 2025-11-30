# font_consistency.py
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
import re, statistics

from app.services.analysis.base_analyzer import BaseAnalyzer
from app.models.analysis_dto import IssueElement, IssueResult, SlideIssueResult
from app.utils.ppt_utils import emu_to_px

# ---------------------------------------------------------
# 🔹 Text Classifier (제목/본문 분류 로직)
# ---------------------------------------------------------
class TextClassifier:
    def __init__(self, slide_height):
        self.slide_height = slide_height

    def classify(self, shape):
        """shape 객체를 받아서 title/body 라벨을 반환"""
        # 1) Placeholder 기반 우선 적용
        try:
            p_type = shape.placeholder_format.type
            if p_type in [
                PP_PLACEHOLDER_TYPE.TITLE,
                PP_PLACEHOLDER_TYPE.CENTER_TITLE,
                PP_PLACEHOLDER_TYPE.SUBTITLE,
            ]:
                return "title"
            elif p_type in [PP_PLACEHOLDER_TYPE.BODY]:
                return "body"
        except ValueError:
            pass

        # 2) 휴리스틱 기반 분류
        top_ratio = shape.top / self.slide_height

        # 첫 번째 run으로 폰트 크기 읽기
        font_size = None
        if shape.has_text_frame and shape.text_frame.paragraphs:
            p = shape.text_frame.paragraphs[0]
            if p.runs and p.runs[0].font.size:
                font_size = p.runs[0].font.size.pt
            else:
                font_size = 20  # 기본 fallback 값
        else:
            return "body"

        text = shape.text.strip()
        word_count = len(text.split())

        score_title = 0
        score_body = 0

        if top_ratio < 0.2:
            score_title += 3
        else:
            score_body += 4

        if font_size >= 36:
            score_title += 2
        else:
            score_body += 2

        if word_count <= 7:
            score_title += 1
        else:
            score_body += 1

        return "title" if score_title >= score_body else "body"


# ---------------------------------------------------------
# 🔹 Font unify checker — unify_test/checker/font_unify_checker.py 복원
# ---------------------------------------------------------
class FontUnifyChecker:
    def __init__(self):
        pass

    def classify_title_shape(self, text: str):
        if not text:
            return "text"

        # 페이지 번호 형태 ex) "1", "2", "10"
        if re.fullmatch(r"[0-9\s\.\-]+", text):
            return "number"

        if len(text) <= 3 and re.search(r"\d", text):
            return "number"

        return "text"

    def get_repr(self, values):
        """최빈값 기반 대표값 산출"""
        if not values:
            return None
        try:
            return round(statistics.mode(values), 1)
        except statistics.StatisticsError:
            return round(statistics.median(values), 1)


# ---------------------------------------------------------
# 🔹 FontConsistencyAnalyzer — 실제 Analyzer 본체
# ---------------------------------------------------------
class FontConsistencyAnalyzer(BaseAnalyzer):
    analyzer_type = "font_consistency"

    def analyze(self, prs):
        slide_height = prs.slide_height
        classifier = TextClassifier(slide_height)
        unify = FontUnifyChecker()

        # 1) 슬라이드별 shape 정보 수집
        parsed = {}  # {slide_idx: [ {shape, label, text, font_size}, ... ]}

        for slide_idx, slide in enumerate(prs.slides):
            shapes_info = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                label = classifier.classify(shape)

                font_size = None
                if shape.text_frame.paragraphs:
                    p = shape.text_frame.paragraphs[0]
                    if p.runs and p.runs[0].font.size:
                        font_size = p.runs[0].font.size.pt

                shapes_info.append(
                    {
                        "shape": shape,
                        "label": label,
                        "text": shape.text.strip(),
                        "font_size": font_size,
                    }
                )

            parsed[slide_idx] = shapes_info

        # 2) 슬라이드별 통계 (number/title 그룹별 대표값)
        slide_stats = {}
        for slide_id, shapes in parsed.items():
            number_sizes, text_sizes = [], []

            for s in shapes:
                if s["label"] != "title":
                    continue
                if not s["font_size"]:
                    continue

                group = unify.classify_title_shape(s["text"])
                if group == "number":
                    number_sizes.append(s["font_size"])
                else:
                    text_sizes.append(s["font_size"])

            avg_number = unify.get_repr(number_sizes)
            avg_text = unify.get_repr(text_sizes)

            slide_stats[slide_id] = {
                "avg_number_size": avg_number,
                "avg_text_size": avg_text,
                # 내부/슬라이드 간 불일치는 여기서만 계산해서,
                # 아래에서 element 이슈 만들 때 details로만 사용
                "number_sizes": number_sizes,
                "text_sizes": text_sizes,
            }

        # 3) 전체 대표값 (모든 슬라이드 title/number 평균의 최빈값/중앙값)
        number_avgs = [
            st["avg_number_size"] for st in slide_stats.values() if st["avg_number_size"]
        ]
        text_avgs = [
            st["avg_text_size"] for st in slide_stats.values() if st["avg_text_size"]
        ]

        global_number_size = unify.get_repr(number_avgs)
        global_text_size = unify.get_repr(text_avgs)

        # 4) 실제 IssueResult: "최빈값에서 벗어나는 개별 title 요소"만 이슈로 생성
        final_results: list[SlideIssueResult] = []

        for slide_id, shapes in parsed.items():
            slide_issues: list[IssueResult] = []
            st = slide_stats[slide_id]

            # 슬라이드 내부 분산 정보
            number_spread = (
                max(st["number_sizes"]) - min(st["number_sizes"])
                if len(st["number_sizes"]) >= 2
                else 0
            )
            text_spread = (
                max(st["text_sizes"]) - min(st["text_sizes"])
                if len(st["text_sizes"]) >= 2
                else 0
            )

            for s in shapes:
                if s["label"] != "title":
                    continue
                if not s["font_size"]:
                    continue

                group = unify.classify_title_shape(s["text"])
                expected = (
                    global_number_size if group == "number" else global_text_size
                )
                if not expected:
                    continue

                # 최빈값에서 벗어난 애들만 이슈로
                if abs(s["font_size"] - expected) <= 1:
                    continue

                shape = s["shape"]
                slide_issues.append(
                    IssueResult(
                        type=self.analyzer_type,
                        message="제목의 글자 크기가 다른 슬라이드와 일관되지 않습니다.",
                        element=IssueElement(
                            shapeId=shape.shape_id,
                            elementIndex=0,
                            bboxLeft=emu_to_px(shape.left),
                            bboxTop=emu_to_px(shape.top),
                            bboxWidth=emu_to_px(shape.width),
                            bboxHeight=emu_to_px(shape.height),
                            text=s["text"],
                            elementType="text",
                        ),
                        details={
                            "currentSize": s["font_size"],
                            "recommendedSize": expected,
                        },
                    )
                )

            if slide_issues:
                final_results.append(
                    SlideIssueResult(slide=slide_id, issues=slide_issues)
                )

        return final_results
