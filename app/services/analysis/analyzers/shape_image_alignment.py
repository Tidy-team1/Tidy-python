from app.services.analysis.base_analyzer import BaseAnalyzer
from app.models.analysis_dto import IssueElement, IssueResult, SlideIssueResult
from app.services.analysis.core.shape_alignment_core import detect_pattern, adjust_shapes_with_pattern
from app.core.logger import logger
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_PX = 9525

ALLOW_TYPES = {
    MSO_SHAPE_TYPE.AUTO_SHAPE,   # 도형
    MSO_SHAPE_TYPE.PICTURE,      # 이미지
    MSO_SHAPE_TYPE.GROUP,        # 그룹
}

# -------------------------------
# width / height 기준으로 그룹화
# -------------------------------
def group_by_size(shapes, ratio=0.1):
    """
    shapes: [{ shapeId, left, top, width, height }]
    width/height가 ±ratio 비율 내에서 비슷한 shape끼리 묶음
    """
    groups = []

    for s in shapes:
        assigned = False
        for g in groups:
            base = g[0]
            if (
                abs(s["width"] - base["width"]) <= base["width"] * ratio and
                abs(s["height"] - base["height"]) <= base["height"] * ratio
            ):
                g.append(s)
                assigned = True
                break
        if not assigned:
            groups.append([s])

    return groups


# -------------------------------
# Analyzer 구현
# -------------------------------
class ShapeImageAlignmentAnalyzer(BaseAnalyzer):
    analyzer_type = "shape_image_alignment"

    def analyze(self, prs):
        slide_results = []

        for slide_idx, slide in enumerate(prs.slides):
            issues = self._analyze_slide(slide_idx, slide)
            slide_results.append(
                SlideIssueResult(slide=slide_idx, issues=issues)
            )

        return slide_results


    def _analyze_slide(self, slide_idx, slide):
        shapes_data = []

        for shape in slide.shapes:

            # ---------- 1) shape_type 필터링 ----------
            st = shape.shape_type
            if st not in ALLOW_TYPES:
                continue

            # ---------- 2) 그룹 안의 서브 도형까지도 포함 ----------
            if st == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    if child.shape_type in ALLOW_TYPES:
                        shapes_data.append({
                            "shapeId": child.shape_id,
                            "left": float(child.left) / EMU_PER_PX,
                            "top": float(child.top) / EMU_PER_PX,
                            "width": float(child.width) / EMU_PER_PX,
                            "height": float(child.height) / EMU_PER_PX,
                        })
                continue  # 그룹 root는 좌표로 정렬하지 않음

            # ---------- 3) 일반 도형/이미지 ----------
            shapes_data.append({
                "shapeId": shape.shape_id,
                "left": float(shape.left) / EMU_PER_PX,
                "top": float(shape.top) / EMU_PER_PX,
                "width": float(shape.width) / EMU_PER_PX,
                "height": float(shape.height) / EMU_PER_PX,
            })

        if not shapes_data:
            return []

        # -------- 1) 크기별 그룹화 --------
        size_groups = group_by_size(shapes_data)

        issues = []

        # -------- 2) 각 그룹 별 패턴 검출 & 조정 --------
        for group in size_groups:
            if len(group) < 2:
                continue

            pattern = detect_pattern(group)
            if pattern == "unknown":
                continue

            adjusted = adjust_shapes_with_pattern(pattern, group)

            # 그룹 전체 bbox 구하기
            bbox = self._group_bbox(group)

            element = IssueElement(
                bboxLeft=bbox["left"],
                bboxTop=bbox["top"],
                bboxWidth=bbox["width"],
                bboxHeight=bbox["height"],
                elementType="shapeGroup"
            )

            details = {
                "pattern": pattern,
                "shapes": [
                    {
                        "shapeId": s["shapeId"],
                        "newLeft": adj["left"],
                        "newTop": adj["top"],
                    }
                    for s, adj in zip(group, adjusted)
                ]
            }

            issues.append(IssueResult(
                type="shape_image_alignment",
                message=f"{pattern} 패턴에서 벗어난 도형들이 있습니다.",
                element=element,
                details=details
            ))

        return issues


    def _group_bbox(self, group):
        lefts = [s["left"] for s in group]
        tops = [s["top"] for s in group]
        rights = [s["left"] + s["width"] for s in group]
        bottoms = [s["top"] + s["height"] for s in group]

        return {
            "left": min(lefts),
            "top": min(tops),
            "width": max(rights) - min(lefts),
            "height": max(bottoms) - min(tops),
        }
