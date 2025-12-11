# app/services/analysis/analyzers/image_contrast.py
import cv2
import numpy as np
from typing import List, Dict, Any
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.models.analysis_dto import IssueElement, IssueResult, SlideIssueResult
from app.services.analysis.base_analyzer import BaseAnalyzer
from app.utils.ppt_utils import emu_to_px

class ImageContrastAnalyzer(BaseAnalyzer):
    analyzer_type = "image_contrast"

    def __init__(self, low_contrast_thresh=40, low_alpha_thresh=180):
        self.low_contrast_thresh = low_contrast_thresh
        self.low_alpha_thresh = low_alpha_thresh

    def analyze(self, prs: Presentation) -> List[SlideIssueResult]:
        results: List[SlideIssueResult] = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_issues: List[IssueResult] = []
            
            for i, shape in enumerate(slide.shapes):
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                
                if not hasattr(shape, "image"):
                    continue
                
                try:
                    image_blob = shape.image.blob
                    nparr = np.frombuffer(image_blob, np.uint8)
                    img = cv2.imdecode(nparr, cv2.IMREAD_UNCHANGED)
                except Exception as e:
                    print(f"[WARN] 이미지 디코딩 실패 (Slide {slide_idx}, Shape {shape.shape_id}): {e}")
                    continue

                # 분석 수행
                analysis = self._analyze_cv2_image(
                    img,
                    low_contrast_thresh=self.low_contrast_thresh,
                    low_alpha_thresh=self.low_alpha_thresh
                )

                if analysis["low_contrast"] or analysis["low_alpha"]:
                    issue_msgs = []
                    
                    # [수정] 조절 값(Adjustment Values) 계산
                    # 백엔드/프론트엔드에서 이 값을 보고 자동으로 수치를 올릴 수 있도록 함
                    recommended: Dict[str, float] = {}
                    
                    if analysis["low_contrast"]:
                        current_c = analysis['contrast_score']
                        required_c = self.low_contrast_thresh
                        diff_c = max(0, required_c - current_c) # 부족한 만큼
                        
                        recommended["contrast_increase"] = round(diff_c, 2) # 올려야 할 대비 값
                        recommended["target_contrast"] = required_c
                        
                        issue_msgs.append(f"대비 낮음({current_c:.1f} < {required_c})")

                    if analysis["low_alpha"]:
                        current_a = analysis['alpha_score']
                        required_a = self.low_alpha_thresh
                        diff_a = max(0, required_a - current_a) # 부족한 만큼
                        
                        recommended["alpha_increase"] = round(diff_a, 2)    # 올려야 할 투명도(알파) 값
                        recommended["target_alpha"] = required_a
                        
                        issue_msgs.append(f"투명도 문제({current_a:.1f} < {required_a})")
                    
                    message = "이미지 시인성 문제: " + ", ".join(issue_msgs)

                    # [수정] details 구성
                    details_payload = {
                        "current": "Image",
                        "contrast": analysis['contrast_score'],
                        "alpha_score": analysis['alpha_score'],
                        
                        # 기준값 정보
                        "required_contrast": self.low_contrast_thresh,
                        "required_alpha": self.low_alpha_thresh,
                        "required": self.low_contrast_thresh, # color_contrast와의 호환성을 위한 대표 required 값
                        
                        # 조절해야 할 수치 정보 전달
                        "recommended": recommended,
                    }

                    slide_issues.append(
                        IssueResult(
                            type=self.analyzer_type,
                            message=message,
                            element=IssueElement(
                                shapeId=shape.shape_id,
                                elementIndex=i,
                                bboxLeft=emu_to_px(shape.left),
                                bboxTop=emu_to_px(shape.top),
                                bboxWidth=emu_to_px(shape.width),
                                bboxHeight=emu_to_px(shape.height),
                                text=None,
                                elementType="image",
                            ),
                            details=details_payload
                        )
                    )

            if slide_issues:
                results.append(SlideIssueResult(slide=slide_idx, issues=slide_issues))

        return results
    
    def _analyze_cv2_image(self, img, low_contrast_thresh, low_alpha_thresh):
        # (기존 로직과 동일)
        if img is None:
            return {
                "status": "error", "contrast_score": 0, "alpha_score": 0,
                "low_contrast": False, "low_alpha": False
            }

        alpha_score = 255.0
        low_alpha = False
        if len(img.shape) == 3 and img.shape[2] == 4:
            alpha_channel = img[:, :, 3]
            alpha_score = float(np.mean(alpha_channel))
            if alpha_score < low_alpha_thresh:
                low_alpha = True

        if len(img.shape) == 3:
            if img.shape[2] == 4:
                gray = cv2.cvtColor(img[:, :, :3], cv2.COLOR_BGR2GRAY)
            else:
                gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        else:
            gray = img 

        contrast_score = float(np.std(gray))
        low_contrast = contrast_score < low_contrast_thresh

        return {
            "status": "ok",
            "contrast_score": round(contrast_score, 2),
            "alpha_score": round(alpha_score, 2),
            "low_contrast": low_contrast,
            "low_alpha": low_alpha,
        }