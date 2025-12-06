# app/services/modify_service.py

import json
import os
import tempfile

from pptx import Presentation
from pptx.util import Pt

from app.core.config import settings
from app.core.logger import logger
from app.utils.s3_utils import download_from_s3, upload_to_s3
from app.services.storage_service import save_file
from app.utils.s3_key_builder import (
    pptx_key_for_version,
    slide_image_key_for_version,
)
from app.services.ppt_to_pdf import convert_ppt_to_pdf
from app.services.pdf_to_images import convert_pdf_to_images

from pptx.enum.shapes import MSO_SHAPE_TYPE

def process_modify(batch):
    """
    전체 수정 작업 메인 엔트리:
    - baseVersion PPTX 불러오기
    - payload.items의 피드백을 차례대로 적용
    - targetVersion PPTX / 이미지 생성 및 업로드
    - 최종 ModifyResult(dict) 반환
    """
    pres_id = batch.presentationId
    space_id = batch.spaceId

    base_v = batch.baseVersion
    target_v = batch.targetVersion

    with tempfile.TemporaryDirectory() as tmpdir:
        base_ppt_path = os.path.join(tmpdir, "base.pptx")
        target_ppt_path = os.path.join(tmpdir, "target.pptx")

        # 1) baseVersion PPT 다운로드
        download_from_s3(
            pptx_key_for_version(space_id, pres_id, base_v),
            base_ppt_path
        )

        prs = Presentation(base_ppt_path)

        # 2) payload.items 반복 → 슬라이드별 수정
        for item in batch.items:
            apply_single_feedback(prs, item)

        # 3) targetVersion PPT 저장
        prs.save(target_ppt_path)

        # 4) target PPT S3 Key
        ppt_key = pptx_key_for_version(space_id, pres_id, target_v)

        # 5) PPT 업로드
        upload_to_s3(
            src_path=target_ppt_path,
            s3_key=ppt_key,
        )

        # 6) PPT → PDF
        pdf_path = convert_ppt_to_pdf(target_ppt_path, tmpdir)

        # 7) PDF → 이미지 바이트 리스트
        image_bytes_list = convert_pdf_to_images(pdf_path)
        slide_count = len(image_bytes_list)

        # 8) 이미지 업로드
        for idx, img_bytes in enumerate(image_bytes_list):
            key = slide_image_key_for_version(space_id, pres_id, target_v, idx)
            save_file(key, img_bytes)

        # 9) slidePrefix
        slide_prefix = f"spaces/{space_id}/presentations/{pres_id}/v{target_v}/slides/"

        # 10) 결과 반환(JSON으로 직렬화됨)
        return {
            "slideCount": slide_count,
            "pptS3Key": ppt_key,
            "slidePrefix": slide_prefix
        }


def apply_single_feedback(prs, item):
    slide = prs.slides[item.slideIndex]

    # ----- detailsJson 파싱 -----
    details_raw = item.detailsJson or "{}"
    details = None

    try:
        # 1차 loads (escape 제거)
        loaded = json.loads(details_raw)

        # loaded가 문자열이면 → 2차 loads 필요
        if isinstance(loaded, str):
            details = json.loads(loaded)
        else:
            details = loaded

    except Exception:
        details = {}
        logger.error(f"[MODIFY] Failed to parse detailsJson: {details_raw}")

    logger.info(
        f"[MODIFY] Apply type='{item.type}', slide={item.slideIndex}, shapeId={item.shapeId}, details={details}"
    )

    # ----- 타입 라우팅 -----
    if item.type == "font_consistency":
        apply_font_consistency(slide, item, details)

    elif item.type == "shape_image_alignment":
        apply_align_shapes(slide, item, details)

    elif item.type == "spelling_grammar":
        apply_text_replacement(slide, item, details)

    else:
        logger.warning(f"[MODIFY] Unknown feedback type: {item.type}")

# -----------------------------
# 유틸 함수
# -----------------------------
def iter_all_shapes(parent):
    """slide or group → 전체 도형 flatten"""
    for shape in parent.shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in iter_all_shapes(shape):
                yield sub



# -----------------------------
# 각 규칙 apply 함수
# -----------------------------
def apply_font_consistency(slide, item, details):
    """
    details 예시:
    {
        "currentSize": 30.19,
        "recommendedSize": 74.0
    }
    """

    if item.shapeId is None:
        logger.warning("[font_consistency] shapeId is None → skip")
        return

    current_size = details.get("currentSize")
    recommended = details.get("recommendedSize")

    if not current_size or not recommended:
        logger.warning(
            f"[font_consistency] currentSize/recommendedSize missing for shapeId={item.shapeId}"
        )
        return

    scale = recommended / current_size

    logger.info(
        f"[font_consistency] shapeId={item.shapeId}, "
        f"currentSize={current_size}, recommendedSize={recommended}, scale={scale:.2f}"
    )

    updated = False

    for shape in slide.shapes:
        if getattr(shape, "shape_id", None) != item.shapeId:
            continue
        if not shape.has_text_frame:
            continue

        # --- 1) 폰트 크기 변경 ---
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font:
                    run.font.size = Pt(recommended)

        # --- 2) 텍스트박스 높이 조정: top 유지 + 아래로 확대 ---
        old_top = shape.top
        old_height = shape.height

        new_height = int(old_height * scale)

        shape.top = old_top            # top 그대로
        shape.height = new_height      # 아래 방향으로 커짐

        logger.info(
            f"[font_consistency] shapeId={item.shapeId} "
            f"height {old_height} → {new_height}, top={shape.top}"
        )

        updated = True

    if not updated:
        logger.warning(f"[font_consistency] No shape updated for shapeId={item.shapeId}")


def apply_align_shapes(slide, item, details):
    shape_positions = details.get("shapes", [])
    if not shape_positions:
        return

    EMU_PER_PX = 9525
    pos_map = {s["shapeId"]: s for s in shape_positions}

    updated = False

    for shape in iter_all_shapes(slide):
        sid = getattr(shape, "shape_id", None)
        if sid not in pos_map:
            continue

        new_left_px = pos_map[sid]["newLeft"]
        new_top_px = pos_map[sid]["newTop"]

        shape.left = int(new_left_px * EMU_PER_PX)
        shape.top = int(new_top_px * EMU_PER_PX)

        updated = True
        logger.info(f"[align_shapes] Updated shapeId={sid}")

    if not updated:
        logger.warning("[align_shapes] No shapes updated")


def apply_text_replacement(slide, item):
    """
    detailsJson 예시:
    {"before": "helo", "after": "hello"}
    """
    if item.shapeId is None:
        return

    before = item.details.get("before")
    after = item.details.get("after")

    if not before or after is None:
        return

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if getattr(shape, "shape_id", None) != item.shapeId:
            continue

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text:
                    run.text = run.text.replace(before, after)
