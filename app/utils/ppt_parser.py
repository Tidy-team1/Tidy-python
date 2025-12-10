# app/utils/ppt_parser.py
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple

from pptx import Presentation
from pptx.enum.dml import MSO_FILL, MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

from app.utils.ppt_utils import emu_to_px

# =========================================================
# 🔹 Internal Utils
# =========================================================

def _rgb_to_hex(rgb: Tuple[int, int, int]) -> Optional[str]:
    """RGB 튜플을 Hex 문자열로 변환"""
    if not rgb:
        return None
    try:
        return "#{:02X}{:02X}{:02X}".format(rgb[0], rgb[1], rgb[2])
    except Exception:
        return None

def _load_theme_colors(pptx_path: Path) -> Dict[str, Tuple[int, int, int]]:
    """PPTX 테마 색상 추출"""
    if not pptx_path:
        return {}
        
    theme_map = {}
    try:
        with zipfile.ZipFile(pptx_path, 'r') as z:
            theme_files = [f for f in z.namelist() if f.startswith("ppt/theme/theme")]
            if not theme_files:
                return {}

            with z.open(theme_files[0]) as theme_xml:
                tree = ET.parse(theme_xml)
                root = tree.getroot()
                ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
                
                theme_elems = root.findall(".//a:themeElements/a:clrScheme/*", ns)
                for elem in theme_elems:
                    name = elem.tag.split("}")[-1].lower()
                    srgb = elem.find(".//a:srgbClr", ns)
                    if srgb is not None and "val" in srgb.attrib:
                        hex_code = srgb.attrib["val"]
                        if len(hex_code) == 6:
                            r, g, b = int(hex_code[0:2], 16), int(hex_code[2:4], 16), int(hex_code[4:6], 16)
                            theme_map[name] = (r, g, b)
                    
                    sys_clr = elem.find(".//a:sysClr", ns)
                    if sys_clr is not None and "lastClr" in sys_clr.attrib:
                         hex_code = sys_clr.attrib["lastClr"]
                         if len(hex_code) == 6:
                            r, g, b = int(hex_code[0:2], 16), int(hex_code[2:4], 16), int(hex_code[4:6], 16)
                            theme_map[name] = (r, g, b)
    except Exception as e:
        print(f"[WARN] Failed to load theme_colors: {e}")
    
    return theme_map

def _get_safe_color(color_obj, theme_map: Dict) -> Optional[str]:
    """색상 객체에서 Hex 코드 추출"""
    if not color_obj:
        return None
    try:
        if color_obj.type == MSO_COLOR_TYPE.RGB:
            return _rgb_to_hex(color_obj.rgb)
        if color_obj.type == MSO_COLOR_TYPE.THEME:
            if color_obj.theme_color:
                theme_key = str(color_obj.theme_color).split('.')[-1].lower()
                rgb_tuple = theme_map.get(theme_key)
                if rgb_tuple:
                    return _rgb_to_hex(rgb_tuple)
    except Exception:
        pass
    return None

def _get_slide_background_color(slide, theme_map: Dict) -> str:
    """슬라이드 배경색 추출"""
    def extract_from_bg(bg):
        if not bg: return None
        fill = bg.fill
        if not fill: return None
        try:
            if fill.type == MSO_FILL.SOLID:
                return _get_safe_color(fill.fore_color, theme_map)
        except:
            pass
        return None

    color = extract_from_bg(slide.background)
    if color: return color
    if slide.slide_layout:
        color = extract_from_bg(slide.slide_layout.background)
        if color: return color
        if slide.slide_layout.slide_master:
            color = extract_from_bg(slide.slide_layout.slide_master.background)
            if color: return color
    return "#FFFFFF"

# =========================================================
# 🔹 Main Logic
# =========================================================

def parse_presentation(prs: Presentation, pptx_path: Path = None) -> Dict:
    """prs 객체를 분석하여 표준 JSON 구조로 변환"""
    theme_map = {}
    if pptx_path and Path(pptx_path).exists():
        theme_map = _load_theme_colors(Path(pptx_path))
    
    slides_data = []
    
    for slide_idx, slide in enumerate(prs.slides):
        bg_hex = _get_slide_background_color(slide, theme_map) or "#FFFFFF"
        elements = []
        palette_colors = set()
        palette_colors.add(bg_hex)
        
        for el_idx, shape in enumerate(slide.shapes):
            elem = {
                "shape_id": shape.shape_id,
                "element_index": el_idx,
                "type": None, # 나중에 결정
                "text": None,
                "left": emu_to_px(shape.left),
                "top": emu_to_px(shape.top),
                "width": emu_to_px(shape.width),
                "height": emu_to_px(shape.height),
                "fontSize": None,
                "isBold": False,
                "textColor": None,
                "fillColor": None,  # [중요] 모든 도형에 대해 추출 시도
                "lineColor": None,
                "image_path": None
            }

            # -------------------------------------------------
            # 1. 색상 추출 (도형/텍스트박스 공통)
            # -------------------------------------------------
            try:
                # Fill Color
                if hasattr(shape, 'fill') and shape.fill.type == MSO_FILL.SOLID:
                    f_color = _get_safe_color(shape.fill.fore_color, theme_map)
                    if f_color:
                        elem["fillColor"] = f_color
                        palette_colors.add(f_color)
                
                # Line Color
                if hasattr(shape, 'line') and shape.line.fill.type == MSO_FILL.SOLID:
                    l_color = _get_safe_color(shape.line.color, theme_map)
                    if l_color:
                        elem["lineColor"] = l_color
                        palette_colors.add(l_color)
            except Exception:
                pass

            # -------------------------------------------------
            # 2. 타입 및 콘텐츠 결정
            # -------------------------------------------------
            
            # (A) 이미지
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                elem["type"] = "image"
                elem["image_path"] = f"extracted_images/slide_{slide_idx+1}_{shape.shape_id}.png"

            # (B) 텍스트가 있는 경우 (텍스트 박스 or 도형 안의 텍스트)
            elif shape.has_text_frame and shape.text.strip():
                elem["type"] = "text"
                elem["text"] = shape.text.strip()
                
                # 텍스트 스타일 추출 (첫 문단 첫 런 기준)
                if shape.text_frame.paragraphs:
                    first_para = shape.text_frame.paragraphs[0]
                    if first_para.runs:
                        run = first_para.runs[0]
                        if run.font.size:
                            elem["fontSize"] = run.font.size.pt
                        elem["isBold"] = run.font.bold or False
                        try:
                            t_color = _get_safe_color(run.font.color, theme_map)
                            if t_color:
                                elem["textColor"] = t_color
                                palette_colors.add(t_color)
                        except:
                            pass

            # (C) 텍스트는 없지만 색상이 있는 경우 -> "도형"으로 취급
            elif elem["fillColor"] is not None or elem["lineColor"] is not None:
                elem["type"] = "shape"
            
            # (D) 그룹 등 기타
            else:
                # 필요한 경우 "group" 등으로 확장 가능하나 일단 넘어감
                pass

            # 타입이 결정된 요소만 추가
            if elem["type"]:
                elements.append(elem)
        
        slides_data.append({
            "slide_index": slide_idx,
            "background_color": bg_hex,
            "palette_hex": list(palette_colors),
            "elements": elements
        })
    
    return {"slides": slides_data}