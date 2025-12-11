import sys
import os
from unittest.mock import MagicMock
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# =========================================================
# 🛑 [중요] image_contrast.py 임포트 전 환경 설정 (Mocking)
# =========================================================
# image_contrast.py가 'app.core.logger' 등을 import 하므로
# 실제 앱 환경이 없어도 에러가 나지 않도록 가짜 모듈을 주입합니다.

mock_logger = MagicMock()
mock_logger.logger = MagicMock()
mock_logger.logger.info = print     # 로그를 print로 대체
mock_logger.logger.warning = print  # 경고를 print로 대체
mock_logger.logger.error = print    # 에러를 print로 대체

sys.modules["app"] = MagicMock()
sys.modules["app.core"] = MagicMock()
sys.modules["app.core.logger"] = mock_logger
sys.modules["app.utils"] = MagicMock()
sys.modules["app.utils.ppt_utils"] = MagicMock()

# =========================================================
# 🚀 모듈 임포트
# =========================================================
# 같은 폴더에 image_contrast.py가 있어야 합니다.
try:
    from image_contrast import apply_image_contrast
except ImportError:
    print("[ERROR] 'image_contrast.py' 파일을 찾을 수 없습니다.")
    print("이 스크립트와 같은 폴더에 image_contrast.py를 위치시켜 주세요.")
    sys.exit(1)

# =========================================================
# ⚙️ 테스트 설정
# =========================================================
INPUT_FILE = r"C:\Users\naeun\capstone_1\before.pptx"
OUTPUT_FILE = r"C:\Users\naeun\capstone_1\after_contrast.pptx"

# DTO 대신 사용할 간단한 클래스
class MockItem:
    def __init__(self, shapeId):
        self.shapeId = shapeId

def run_test():
    print(f"[INIT] PPT 파일 로드 중: {INPUT_FILE}")
    
    if not os.path.exists(INPUT_FILE):
        print(f"[ERROR] 입력 파일이 존재하지 않습니다: {INPUT_FILE}")
        return

    prs = Presentation(INPUT_FILE)
    
    # 첫 번째 슬라이드 사용
    if not prs.slides:
        print("[ERROR] 슬라이드가 없습니다.")
        return
    
    slide = prs.slides[0]
    
    # 슬라이드에서 첫 번째 이미지 찾기
    target_shape = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            target_shape = shape
            break
            
    if not target_shape:
        print("[ERROR] 첫 번째 슬라이드에서 이미지를 찾을 수 없습니다.")
        return

    print(f"[INFO] 대상 이미지 발견 (Shape ID: {target_shape.shape_id})")

    # ---------------------------------------------------------
    # 🎯 [시나리오] 흐릿한 이미지를 선명하게 만들기
    # ---------------------------------------------------------
    # 현재 대비(contrast)는 낮고(예: 10.0), 목표(target)는 높게(예: 60.0) 설정
    details = {
        "current": "Image",
        "contrast": 9.89,       # (가정) 현재 대비가 매우 낮음
        "recommended": {
            "target_contrast": 60.0,    # 목표 대비 (높을수록 대비가 강해짐)
            "contrast_increase": 50.0,  # 차이값
            
            # 투명도 조절이 필요 없다면 아래 두 줄은 주석 처리 가능
            # "target_alpha": 255.0, 
            # "alpha_increase": 0.0
        }
    }

    item = MockItem(shapeId=target_shape.shape_id)

    print("[ACTION] 이미지 보정(Contrast Enhancement) 적용 중...")
    
    # 🔥 함수 실행
    apply_image_contrast(slide, item, details)

    # 결과 저장
    prs.save(OUTPUT_FILE)
    print("\n" + "="*50)
    print(f"✅ 테스트 완료! 결과 파일이 생성되었습니다.")
    print(f"📂 저장 경로: {OUTPUT_FILE}")
    print("="*50)

if __name__ == "__main__":
    run_test()